from pptx import Presentation


class PptxGenerator:
    @staticmethod
    def generate(content, filename="example", **kwargs):
        # Create a presentation object
        prs = Presentation()

        # Add a slide with title and subtitle layout
        slide_layout = prs.slide_layouts[0]  # 0 is the layout index for title slide
        slide = prs.slides.add_slide(slide_layout)

        # Add title and subtitle to the slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Hello, World!"
        subtitle.text = content

        # Add a slide with title and content layout
        slide_layout = prs.slide_layouts[
            1
        ]  # 1 is the layout index for title and content slide
        slide = prs.slides.add_slide(slide_layout)

        # Add title and content to the slide
        title = slide.shapes.title
        title.text = "Adding a Bullet Slide"

        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "This is a bullet slide."

        # Add bullets to the text frame
        p = text_frame.add_paragraph()
        p.text = "First bullet"
        p.level = 0

        p = text_frame.add_paragraph()
        p.text = "Second bullet"
        p.level = 1

        p = text_frame.add_paragraph()
        p.text = "Third bullet"
        p.level = 2

        # Add a slide with an image
        slide_layout = prs.slide_layouts[5]  # 5 is the layout index for a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add an image to the slide
        """
        from pptx.util import Inches
        img_path = 'path/to/your/image.png'
        left = Inches(1)
        top = Inches(1)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)
        
        """

        # Save the presentation
        prs.save(f"{filename}.pptx")
